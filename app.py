import streamlit as st
import pandas as pd

import io
import os
import re
import hashlib
from datetime import datetime, timezone

import gspread
from google.oauth2.service_account import Credentials

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

# =======================
# Config / Constants
# =======================
APP_VERSION = "SLEI-v2.0-pilot"
TEMPLATE_PATH = "SLEI_Dashboard_TEMPLATE.pptx"  # tokenized pptx in repo root

ITEMS = [
    (1, "Pause long enough to respond thoughtfully in high-pressure situations", "Sight"),
    (2, "Take deliberate actions to influence long-term outcomes", "Sight"),
    (3, "Test your assumptions before responding to a situation", "Sight"),
    (4, "Recognize predictable patterns in your reactions and adjust your response accordingly", "Sight"),
    (5, "Align major commitments with what you identify as most important in your role", "Tenacity"),
    (6, "Adjust workload or boundaries to sustain performance", "Tenacity"),
    (7, "Remove or delegate commitments that do not require your direct involvement", "Tenacity"),
    (8, "Intentionally focus your time on the highest-level responsibilities your role is designed to perform", "Tenacity"),
    (9, "Take timely, considered action on important issues even when outcomes are uncertain", "Ability"),
    (10, "When challenges arise, focus first on actions within your control", "Ability"),
    (11, "Renegotiate commitments before they become risks", "Ability"),
    (12, "Intentionally develop others’ capability so they can operate with greater ownership and independence", "Ability"),
    (13, "Before launching work, ensure clarity about outcomes, ownership, resources, and how the team will operate", "Results"),
    (14, "Hold check-ins to ensure expectations remain clear", "Results"),
    (15, "Delegate work with clear expectations rather than retaining it", "Results"),
    (16, "Effectively influence stakeholders beyond your formal authority to move important work forward", "Results"),
]

FREQ_OPTIONS = ["Rarely", "Occasionally", "Sometimes", "Often", "Consistently", "Not applicable to my role"]
CHANGE_OPTIONS = ["Much less often", "Slightly less often", "About the same", "Slightly more often", "Much more often"]

FREQ_MAP = {
    "Rarely": 1,
    "Occasionally": 2,
    "Sometimes": 3,
    "Often": 4,
    "Consistently": 5,
    "Not applicable to my role": None,
}
CHANGE_MAP = {
    "Much less often": -2,
    "Slightly less often": -1,
    "About the same": 0,
    "Slightly more often": 1,
    "Much more often": 2,
}

# Dashboard token keys (must match {{TOKEN}} in template)
TOK = {
    "TITLE": "{{DASHBOARD_TITLE}}",
    "SUBTITLE": "{{DASHBOARD_SUBTITLE}}",
    "INTRO": "{{DASHBOARD_INTRO}}",

    "CONS_SCORE": "{{CONSISTENCY_SCORE}}",
    "CONS_BAND": "{{CONSISTENCY_BAND}}",
    "CONS_INTERP": "{{CONSISTENCY_INTERPRETATION}}",

    "GROWTH_AVG": "{{GROWTH_AVG}}",
    "GROWTH_SUM": "{{GROWTH_SUMMARY}}",
    "GROWTH_INTERP": "{{GROWTH_INTERPRETATION}}",

    "SIGHT": "{{SIGHT_SCORE}}",
    "TENACITY": "{{TENACITY_SCORE}}",
    "ABILITY": "{{ABILITY_SCORE}}",
    "RESULTS": "{{RESULTS_SCORE}}",

    "STRONGEST_GROWTH_DOM": "{{STRONGEST_GROWTH_DOMAIN}}",
    "STRONGEST_GROWTH_DETAIL": "{{STRONGEST_GROWTH_DETAIL}}",

    "TOP_GROWTH_1": "{{TOP_GROWTH_ITEM_1}}",
    "TOP_GROWTH_2": "{{TOP_GROWTH_ITEM_2}}",
    "TOP_GROWTH_WHY": "{{TOP_GROWTH_WHY}}",
    "TOP_GROWTH_PROMPT": "{{TOP_GROWTH_PROMPT}}",

    "TOP_OPP_1": "{{TOP_OPPORTUNITY_ITEM_1}}",
    "TOP_OPP_2": "{{TOP_OPPORTUNITY_ITEM_2}}",
    "TOP_OPP_NEXT": "{{TOP_OPPORTUNITY_NEXTSTEP}}",

    "ACTION_PROMPT": "{{ACTIONPLAN_PROMPT}}",
}

# =======================
# Helpers
# =======================

def safe_mean(vals):
    vals = [v for v in vals if isinstance(v, (int, float))]
    return sum(vals) / len(vals) if vals else None


def round1(x):
    return None if x is None else round(x, 1)


def overall_descriptor(score):
    if score is None:
        return "Not scored"
    if score >= 4.5:
        return "Consistently"
    if score >= 4.0:
        return "Often"
    if score >= 3.0:
        return "Sometimes"
    if score >= 2.0:
        return "Occasionally"
    return "Rarely"


def consistency_interpretation(score):
    # Short, non-judgy; optimized to fit box.
    if score is None:
        return "No overall score (all items marked Not applicable)."
    band = overall_descriptor(score)
    if band == "Consistently":
        return "Strong habit strength — these behaviors show up reliably, even under pressure."
    if band == "Often":
        return "Solid consistency — these behaviors show up most of the time; keep building repeatable habits."
    if band == "Sometimes":
        return "Good foundation — these behaviors show up, with room to make them more consistent in real situations."
    if band == "Occasionally":
        return "Early momentum — choose 1–2 behaviors to practice intentionally and build a reliable cadence."
    return "Starting point — pick one behavior, practice it in a real situation, and keep it small and repeatable."


def growth_interpretation(avg_change):
    if avg_change is None:
        return "No change score available."
    if avg_change >= 1.25:
        return "Large positive shift in how often you’re applying the behaviors."
    if avg_change >= 0.5:
        return "Clear positive shift in how often you’re applying the behaviors."
    if avg_change > 0:
        return "Modest positive shift in how often you’re applying the behaviors."
    if avg_change == 0:
        return "No overall change in reported application (some items may still have improved)."
    if avg_change <= -1.25:
        return "Large decrease in reported application (consider context, workload, or role changes)."
    if avg_change <= -0.5:
        return "Decrease in reported application (consider context, workload, or role changes)."
    return "Slight decrease in reported application (consider context, workload, or role changes)."


def format_bullets(q_items):
    """q_items: list[(qid, text)]"""
    lines = []
    for qid, text in q_items:
        lines.append(f"• Q{qid}: {text}")
    return "
".join(lines)


def open_sheet():
    creds_info = st.secrets["gcp_service_account"]
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]
    creds = Credentials.from_service_account_info(creds_info, scopes=scopes)
    gc = gspread.authorize(creds)
    return gc.open(st.secrets["sheet_name"]).worksheet(st.secrets["worksheet_name"])


def append_row_to_sheet(ws, row):
    ws.append_row(row, value_input_option="RAW")


def init_state():
    st.session_state.setdefault("step", 1)

    # Context
    st.session_state.setdefault("role_anchor", None)
    st.session_state.setdefault("role_anchor_other", "")
    st.session_state.setdefault("profession", None)
    st.session_state.setdefault("profession_other", "")
    st.session_state.setdefault("years", None)
    st.session_state.setdefault("scope", None)
    st.session_state.setdefault("scope_other", "")

    # Responses
    st.session_state.setdefault("freq_sel", {qid: None for qid, _, _ in ITEMS})
    st.session_state.setdefault("chg_sel", {qid: None for qid, _, _ in ITEMS})

    # Feedback / optional contact
    st.session_state.setdefault("improve_feedback", "")
    st.session_state.setdefault("testimonial", "")
    st.session_state.setdefault("willing_contact", False)
    st.session_state.setdefault("contact_name", "")
    st.session_state.setdefault("contact_email", "")

    # Dashboard (per-session)
    st.session_state.setdefault("dashboard_bytes", None)
    st.session_state.setdefault("dashboard_filename", None)
    st.session_state.setdefault("dashboard_generated_at", None)
    st.session_state.setdefault("last_submit_fingerprint", None)


def go_next():
    st.session_state.step += 1


def go_prev():
    st.session_state.step -= 1


def compute_scores():
    freq_num = {qid: FREQ_MAP.get(st.session_state.freq_sel.get(qid)) for qid, _, _ in ITEMS}
    chg_num = {qid: CHANGE_MAP.get(st.session_state.chg_sel.get(qid)) for qid, _, _ in ITEMS}

    # Overall (current frequency)
    freq_vals = [v for v in freq_num.values() if isinstance(v, (int, float))]
    overall = round1(safe_mean(freq_vals))
    overall_desc = overall_descriptor(overall)

    # Domain averages (current frequency)
    domains = sorted({d for _, _, d in ITEMS})
    domain_scores = {}
    for dom in domains:
        dom_vals = [freq_num[qid] for qid, _, d in ITEMS if d == dom and isinstance(freq_num.get(qid), (int, float))]
        domain_scores[dom] = round1(safe_mean(dom_vals))

    # Growth avg (only for items not NA)
    non_na_ids = [qid for qid, _, _ in ITEMS if freq_num.get(qid) is not None]
    chg_vals = [chg_num[qid] for qid in non_na_ids if isinstance(chg_num.get(qid), (int, float))]
    avg_change = safe_mean(chg_vals)
    growth = (avg_change is not None) and (avg_change > 0)

    # Growth count
    inc_count = sum(1 for v in chg_vals if isinstance(v, (int, float)) and v > 0)

    return freq_num, chg_num, overall, overall_desc, domain_scores, avg_change, growth, inc_count, non_na_ids


def required_missing_step1():
    missing = []

    ra = st.session_state.role_anchor
    if not ra:
        missing.append("Leadership role")
    elif ra == "Other" and not st.session_state.role_anchor_other.strip():
        missing.append("Leadership role (Other)")

    prof = st.session_state.profession
    if not prof:
        missing.append("Profession type")
    elif prof == "Other" and not st.session_state.profession_other.strip():
        missing.append("Profession type (Other)")

    if not st.session_state.years:
        missing.append("Years of experience")

    sc = st.session_state.scope
    if not sc:
        missing.append("Leadership scope")
    elif sc == "Other" and not st.session_state.scope_other.strip():
        missing.append("Leadership scope (Other)")

    return missing


def required_missing_freq():
    return [qid for qid, _, _ in ITEMS if st.session_state.freq_sel.get(qid) is None]


def required_missing_change(non_na_ids):
    return [qid for qid in non_na_ids if st.session_state.chg_sel.get(qid) is None]


def build_fingerprint():
    """A stable-ish fingerprint to avoid generating multiple dashboards for the same submission in-session."""
    parts = [
        str(st.session_state.role_anchor or ""),
        st.session_state.role_anchor_other.strip(),
        str(st.session_state.profession or ""),
        st.session_state.profession_other.strip(),
        str(st.session_state.years or ""),
        str(st.session_state.scope or ""),
        st.session_state.scope_other.strip(),
    ]
    # Add responses
    for qid, _, _ in ITEMS:
        parts.append(str(st.session_state.freq_sel.get(qid) or ""))
    for qid, _, _ in ITEMS:
        parts.append(str(st.session_state.chg_sel.get(qid) or ""))
    raw = "|".join(parts)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:16]


# =======================
# PPTX token replacement
# =======================

def _set_autofit(shape):
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def replace_tokens_in_ppt(prs: Presentation, token_map: dict) -> None:
    """Replace {{TOKEN}} occurrences across all runs.

    Important: we do NOT touch colors; prior attempts to read/clone colors can crash
    on theme-based fonts ("_NoneColor").
    """
    # Token-specific sizing (to avoid spillover in tight boxes)
    small_tokens = {
        TOK["TOP_GROWTH_WHY"],
        TOK["TOP_GROWTH_PROMPT"],
        TOK["TOP_OPP_NEXT"],
        TOK["CONS_INTERP"],
        TOK["GROWTH_INTERP"],
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            touched = False
            for p in tf.paragraphs:
                for r in p.runs:
                    if not r.text:
                        continue
                    for k, v in token_map.items():
                        if k in r.text:
                            r.text = r.text.replace(k, v)
                            touched = True
                            # conservative font sizing for dense blocks
                            if k in small_tokens:
                                try:
                                    r.font.size = Pt(11)
                                except Exception:
                                    pass
                            elif k in {TOK["TOP_GROWTH_1"], TOK["TOP_GROWTH_2"], TOK["TOP_OPP_1"], TOK["TOP_OPP_2"]}:
                                try:
                                    r.font.size = Pt(14)
                                except Exception:
                                    pass
            if touched:
                _set_autofit(shape)


def build_dashboard_pptx_bytes(token_map: dict, template_path: str) -> bytes:
    prs = Presentation(template_path)
    replace_tokens_in_ppt(prs, token_map)
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


# =======================
# UI styling
# =======================

SURVEY_CSS = """
<style>
/* Make pages feel less airy + improve readability */
.block-container { padding-top: 1.2rem; padding-bottom: 2.0rem; }

/* Headings */
h1, h2, h3 { letter-spacing: -0.02em; }

/* Radio groups: reduce vertical space */
div[data-baseweb="radio"] > div { gap: 0.75rem !important; }

/* Question text */
.slei-q { font-size: 1.02rem; font-weight: 600; margin: 1.15rem 0 0.35rem 0; }

/* Slightly smaller radio labels */
div[data-baseweb="radio"] label { font-size: 0.95rem; }

/* Compact captions */
.slei-cap { color: rgba(49,51,63,0.72); margin-top: -0.3rem; margin-bottom: 0.75rem; }

/* Section separators */
.slei-divider { border-top: 1px solid rgba(49,51,63,0.14); margin: 1.1rem 0; }
</style>
"""


def q_header(text: str):
    st.markdown(f"<div class='slei-q'>{text}</div>", unsafe_allow_html=True)


def divider():
    st.markdown("<div class='slei-divider'></div>", unsafe_allow_html=True)


# =======================
# App
# =======================

init_state()

st.set_page_config(page_title="SLEI v2.0", layout="wide")
st.markdown(SURVEY_CSS, unsafe_allow_html=True)

st.title("STAR Leadership Effectiveness Index (SLEI) – v2 Pilot")

st.markdown(
    """**Purpose**  
This assessment supports your continued development and helps us improve the program.

**Structure**  
You’ll answer 16 questions about key leadership behaviors in two sections:
1) **Current Frequency** — how often you apply each behavior now (after the course)
2) **Change in Frequency** — how your current application compares to before the course

Because leadership looks different across contexts, select one role and use it consistently so your responses are accurate and comparable.
"""
)

# -----------------------
# Step 1 of 5 — Context
# -----------------------
if st.session_state.step == 1:
    st.header("Step 1 of 5 — Context")

    st.session_state.role_anchor = st.selectbox(
        "Which single leadership role will you use as your reference point? (required)",
        [
            "My primary professional/employer role",
            "A volunteer or board leadership role",
            "A family or community leadership role",
            "Other",
        ],
        index=None,
        placeholder="Select one…",
    )
    if st.session_state.role_anchor == "Other":
        st.session_state.role_anchor_other = st.text_input(
            "If Other, please specify (required)",
            value=st.session_state.role_anchor_other,
        )

    st.session_state.profession = st.selectbox(
        "Profession type (required)",
        ["Student", "Resident", "Pharmacy Technician", "Pharmacist", "Other"],
        index=None,
        placeholder="Select one…",
    )
    if st.session_state.profession == "Other":
        st.session_state.profession_other = st.text_input(
            "If Other, please specify (required)",
            value=st.session_state.profession_other,
        )

    st.session_state.years = st.selectbox(
        "Years of experience (required)",
        ["0–2", "3–5", "6–10", "11–15", "16–20", "21+"],
        index=None,
        placeholder="Select one…",
    )

    st.session_state.scope = st.selectbox(
        "Leadership scope (required)",
        [
            "Individual contributor (no direct reports)",
            "Supervisor / Manager of individuals",
            "Manager of managers",
            "Senior leader / executive",
            "Other",
        ],
        index=None,
        placeholder="Select one…",
    )
    if st.session_state.scope == "Other":
        st.session_state.scope_other = st.text_input(
            "If Other, please specify (required)",
            value=st.session_state.scope_other,
        )

    cols = st.columns([1, 7])
    with cols[0]:
        st.button("Next →", type="primary", on_click=go_next)

    missing = required_missing_step1()
    if missing:
        st.info("To continue, complete: " + ", ".join(missing))

    if st.session_state.step == 2 and missing:
        st.session_state.step = 1


# -----------------------
# Step 2 of 5 — Current frequency
# -----------------------
elif st.session_state.step == 2:
    st.header("Step 2 of 5 — Current frequency")
    st.markdown(
        "<div class='slei-cap'>For each item, select how often you apply it now. Choose ‘Not applicable’ if it does not apply to the role you selected.</div>",
        unsafe_allow_html=True,
    )

    for qid, text, _ in ITEMS:
        q_header(f"Q{qid}. {text}")
        st.session_state.freq_sel[qid] = st.radio(
            label=f"freq_{qid}",
            options=FREQ_OPTIONS,
            index=None,
            horizontal=True,
            label_visibility="collapsed",
            key=f"freq_{qid}",
        )

    divider()

    cols = st.columns([1, 1, 6])
    with cols[0]:
        st.button("← Back", on_click=go_prev)
    with cols[1]:
        st.button("Next →", type="primary", on_click=go_next)

    missing_qs = required_missing_freq()
    if missing_qs:
        st.info(f"To continue, answer all current-frequency items (missing: {len(missing_qs)}).")

    if st.session_state.step == 3 and missing_qs:
        st.session_state.step = 2


# -----------------------
# Step 3 of 5 — Change in frequency
# -----------------------
elif st.session_state.step == 3:
    st.header("Step 3 of 5 — Change in frequency")
    st.markdown(
        "<div class='slei-cap'>You’ll only see change questions for items you did not mark as ‘Not applicable.’</div>",
        unsafe_allow_html=True,
    )

    freq_num, chg_num, overall, overall_desc, domain_scores, avg_change, growth, inc_count, non_na_ids = compute_scores()

    if not non_na_ids:
        st.warning("You marked all items as ‘Not applicable.’ You can still provide optional feedback.")
    else:
        for qid, text, _ in ITEMS:
            if qid not in non_na_ids:
                continue
            q_header(f"Q{qid}. {text}")
            st.session_state.chg_sel[qid] = st.radio(
                label=f"chg_{qid}",
                options=CHANGE_OPTIONS,
                index=None,
                horizontal=True,
                label_visibility="collapsed",
                key=f"chg_{qid}",
            )

    divider()

    cols = st.columns([1, 1, 6])
    with cols[0]:
        st.button("← Back", on_click=go_prev)
    with cols[1]:
        st.button("Next →", type="primary", on_click=go_next)

    missing_chg = required_missing_change(non_na_ids)
    if non_na_ids and missing_chg:
        st.info(f"To continue, answer all change items shown (missing: {len(missing_chg)}).")

    if st.session_state.step == 4 and non_na_ids and missing_chg:
        st.session_state.step = 3


# -----------------------
# Step 4 of 5 — Optional feedback
# -----------------------
elif st.session_state.step == 4:
    st.header("Step 4 of 5 — Optional feedback")

    st.session_state.improve_feedback = st.text_area(
        "Any suggestions to improve the course structure, processes, systems, or curriculum? (optional)",
        value=st.session_state.improve_feedback,
        height=160,
    )

    divider()

    st.subheader("Testimonial (optional)")
    st.caption(
        "If you choose to share a testimonial, it helps others understand the value of the program. "
        "A helpful testimonial often includes: (1) what changed for you, (2) one concrete example, and (3) what you’d say to someone considering the program."
    )

    st.session_state.testimonial = st.text_area(
        "If you’d like, share a short testimonial or comment about the program (optional)",
        value=st.session_state.testimonial,
        height=160,
    )

    st.session_state.willing_contact = st.checkbox(
        "I’m open to being contacted about using my testimonial (optional)",
        value=st.session_state.willing_contact,
    )

    cols = st.columns([1, 1, 6])
    with cols[0]:
        st.button("← Back", on_click=go_prev)
    with cols[1]:
        st.button("Next →", type="primary", on_click=go_next)


# -----------------------
# Step 5 of 5 — Review / Submit + auto-dashboard
# -----------------------
elif st.session_state.step == 5:
    st.header("Step 5 of 5 — Submit")

    if st.session_state.willing_contact:
        st.caption("If you’re comfortable, provide contact details so we can follow up about using your testimonial.")
        st.session_state.contact_name = st.text_input("Name (optional)", value=st.session_state.contact_name)
        st.session_state.contact_email = st.text_input("Email (optional)", value=st.session_state.contact_email)

    divider()

    cols = st.columns([1, 1, 6])
    with cols[0]:
        st.button("← Back", on_click=go_prev)
    with cols[1]:
        submitted = st.button("Submit", type="primary")

    if submitted:
        # Final validation
        missing = required_missing_step1()
        if missing:
            st.error("Missing required fields: " + ", ".join(missing))
            st.stop()

        missing_freq = required_missing_freq()
        if missing_freq:
            st.error("Missing current-frequency answers.")
            st.stop()

        freq_num, chg_num, overall, overall_desc, domain_scores, avg_change, growth, inc_count, non_na_ids = compute_scores()
        missing_chg = required_missing_change(non_na_ids)
        if non_na_ids and missing_chg:
            st.error("Missing change answers for one or more items.")
            st.stop()

        # Resolve 'Other' fields
        role_anchor = st.session_state.role_anchor
        if role_anchor == "Other":
            role_anchor = st.session_state.role_anchor_other.strip()

        profession = st.session_state.profession
        if profession == "Other":
            profession = st.session_state.profession_other.strip()

        scope = st.session_state.scope
        if scope == "Other":
            scope = st.session_state.scope_other.strip()

        # Prevent duplicate generation in the same session
        fp = build_fingerprint()
        if st.session_state.last_submit_fingerprint == fp and st.session_state.dashboard_bytes is not None:
            st.info("Dashboard already generated for this submission. Use the download button below.")
        else:
            # ----- Build dashboard token map -----
            # Pick strongest reported growth items (largest positive change)
            growth_candidates = []
            for qid, text, dom in ITEMS:
                if freq_num.get(qid) is None:
                    continue
                cv = chg_num.get(qid)
                if isinstance(cv, (int, float)):
                    growth_candidates.append((cv, qid, text, dom))
            growth_candidates.sort(key=lambda t: (t[0], -t[1]), reverse=True)
            top_growth = [(qid, text) for _v, qid, text, _d in growth_candidates if _v > 0][:2]

            # Opportunities: lowest current-frequency scores (excluding N/A)
            opp_candidates = []
            for qid, text, dom in ITEMS:
                fv = freq_num.get(qid)
                if isinstance(fv, (int, float)):
                    opp_candidates.append((fv, qid, text, dom))
            opp_candidates.sort(key=lambda t: (t[0], t[1]))
            top_opp = [(qid, text) for _v, qid, text, _d in opp_candidates][:2]

            # Domain with strongest growth (avg change within domain)
            dom_growth = {}
            for dom in sorted({d for _, _, d in ITEMS}):
                dom_vals = []
                for qid, _t, d in ITEMS:
                    if d != dom:
                        continue
                    if freq_num.get(qid) is None:
                        continue
                    v = chg_num.get(qid)
                    if isinstance(v, (int, float)):
                        dom_vals.append(v)
                dom_growth[dom] = safe_mean(dom_vals)

            best_growth_dom = None
            best_growth_val = None
            for dom, val in dom_growth.items():
                if val is None:
                    continue
                if best_growth_val is None or val > best_growth_val:
                    best_growth_val = val
                    best_growth_dom = dom

            # Short “detail” line
            if best_growth_dom is None:
                best_growth_dom = "—"
                best_growth_detail = "No change data available"
            else:
                best_growth_detail = f"Most gains in {best_growth_dom.lower()} behaviors"

            # Domain scores (current)
            sight = domain_scores.get("Sight")
            tenacity = domain_scores.get("Tenacity")
            ability = domain_scores.get("Ability")
            results = domain_scores.get("Results")

            # Growth avg display (+x.x avg)
            avg_change_r1 = round1(avg_change)
            if avg_change_r1 is None:
                growth_avg_disp = "—"
                growth_sum = "No change data"
            else:
                sign = "+" if avg_change_r1 > 0 else ""
                growth_avg_disp = f"{sign}{avg_change_r1:.1f} avg"
                growth_sum = f"Increased in {inc_count} of {len(non_na_ids)} behaviors"

            cons_score_disp = "—" if overall is None else f"{overall:.1f} / 5"
            cons_band_disp = overall_desc
            cons_interp = consistency_interpretation(overall)

            intro = (
                "This dashboard is not a course grade. It summarizes your self-reported application of key leadership behaviors "
                "and the change you reported since before the course."
            )

            # Keep these tight to avoid text collisions
            top_growth_why = (
                "Why this matters: small shifts in application compound — especially when you build repeatable habits in real situations."
            )
            top_growth_prompt = (
                "Use this as a prompt: pick one real situation in the next 60–90 days where you’ll practice intentionally."
            )
            opp_next = "Identify one upcoming situation to practice each behavior intentionally."

            action_prompt = (
                "Complete this 60–90 day action plan (progress > perfection):

"
                "Primary behavior to strengthen (pick ONE): ________________________________

"
                "Situation to practice (next 60–90 days): _________________________________

"
                "What will success look like (observable): _________________________________

"
                "First next step (within 7 days): ________________________________________"
            )

            token_map = {
                TOK["TITLE"]: "SLEI Participant Dashboard",
                TOK["SUBTITLE"]: f"STAR Leadership Effectiveness Index • {APP_VERSION}",
                TOK["INTRO"]: intro,

                TOK["CONS_SCORE"]: cons_score_disp,
                TOK["CONS_BAND"]: cons_band_disp,
                TOK["CONS_INTERP"]: f"Overall: {overall if overall is not None else '—'}/5 ({overall_desc}) — {cons_interp}",

                TOK["GROWTH_AVG"]: growth_avg_disp,
                TOK["GROWTH_SUM"]: growth_sum,
                TOK["GROWTH_INTERP"]: growth_interpretation(avg_change_r1),

                TOK["SIGHT"]: "—" if sight is None else f"{sight:.1f}",
                TOK["TENACITY"]: "—" if tenacity is None else f"{tenacity:.1f}",
                TOK["ABILITY"]: "—" if ability is None else f"{ability:.1f}",
                TOK["RESULTS"]: "—" if results is None else f"{results:.1f}",

                TOK["STRONGEST_GROWTH_DOM"]: best_growth_dom,
                TOK["STRONGEST_GROWTH_DETAIL"]: best_growth_detail,

                TOK["TOP_GROWTH_1"]: top_growth[0][0] if False else "",  # placeholder; overwritten below
                TOK["TOP_GROWTH_2"]: top_growth[0][0] if False else "",  # placeholder; overwritten below
                TOK["TOP_GROWTH_WHY"]: top_growth_why,
                TOK["TOP_GROWTH_PROMPT"]: top_growth_prompt,

                TOK["TOP_OPP_1"]: "",
                TOK["TOP_OPP_2"]: "",
                TOK["TOP_OPP_NEXT"]: opp_next,

                TOK["ACTION_PROMPT"]: action_prompt,
            }

            # Populate top-growth / opp items (keep lines short; template already includes bullets)
            token_map[TOK["TOP_GROWTH_1"]] = top_growth[0][0] if False else ""  # no-op
            # Real: the placeholders are whole bullet lines in template, so we supply full text without extra bullets.
            token_map[TOK["TOP_GROWTH_1"]] = top_growth[0][0] if False else ""  # overwritten next
            token_map[TOK["TOP_GROWTH_2"]] = top_growth[0][0] if False else ""  # overwritten next

            if len(top_growth) >= 1:
                token_map[TOK["TOP_GROWTH_1"]] = f"Q{top_growth[0][0]}: {top_growth[0][1]}" if isinstance(top_growth[0][0], int) else ""
            else:
                token_map[TOK["TOP_GROWTH_1"]] = ""
            if len(top_growth) >= 2:
                token_map[TOK["TOP_GROWTH_2"]] = f"Q{top_growth[1][0]}: {top_growth[1][1]}" if isinstance(top_growth[1][0], int) else ""
            else:
                token_map[TOK["TOP_GROWTH_2"]] = ""

            if len(top_opp) >= 1:
                token_map[TOK["TOP_OPP_1"]] = f"Q{top_opp[0][0]}: {top_opp[0][1]}" if isinstance(top_opp[0][0], int) else ""
            else:
                token_map[TOK["TOP_OPP_1"]] = ""
            if len(top_opp) >= 2:
                token_map[TOK["TOP_OPP_2"]] = f"Q{top_opp[1][0]}: {top_opp[1][1]}" if isinstance(top_opp[1][0], int) else ""
            else:
                token_map[TOK["TOP_OPP_2"]] = ""

            # ----- Save response to Google Sheets (with dashboard fields) -----
            now_utc = datetime.now(timezone.utc).isoformat()
            dash_generated_at = now_utc
            dash_filename = f"SLEI_Dashboard_{dash_generated_at[:19].replace(':','-')}_{fp}.pptx"

            # Build final row
            row = [
                now_utc,
                APP_VERSION,
                role_anchor,
                profession,
                st.session_state.years,
                scope,
                str(overall) if overall is not None else "",
                overall_desc,
                str(round1(avg_change_r1)) if avg_change_r1 is not None else "",
                "Yes" if growth else "No",
                "Yes" if st.session_state.willing_contact else "No",
                (st.session_state.testimonial or "").strip(),
                (st.session_state.improve_feedback or "").strip(),
                (st.session_state.contact_name or "").strip(),
                (st.session_state.contact_email or "").strip(),
            ]

            # freq_1..16
            for qid in range(1, 17):
                v = freq_num.get(qid)
                row.append(str(v) if isinstance(v, (int, float)) else "")

            # chg_1..16 (blank for NA)
            for qid in range(1, 17):
                v = chg_num.get(qid)
                row.append(str(v) if (qid in non_na_ids and isinstance(v, (int, float))) else "")

            # dashboard fields
            row.append(dash_generated_at)
            row.append(dash_filename)

            try:
                ws = open_sheet()
                append_row_to_sheet(ws, row)
            except Exception as e:
                st.error("Could not save to Google Sheets (secrets / sheet setup may be missing or incomplete).")
                st.exception(e)
                st.stop()

            # ----- Build PPTX bytes & store in session -----
            try:
                ppt_bytes = build_dashboard_pptx_bytes(token_map=token_map, template_path=TEMPLATE_PATH)
                st.session_state.dashboard_bytes = ppt_bytes
                st.session_state.dashboard_filename = dash_filename
                st.session_state.dashboard_generated_at = dash_generated_at
                st.session_state.last_submit_fingerprint = fp
            except Exception as e:
                st.error("Dashboard generation failed.")
                st.exception(e)
                st.stop()

        # Download UI
        if st.session_state.dashboard_bytes is not None:
            st.success("Submitted. Your dashboard is ready to download.")
            st.download_button(
                "Download dashboard (PPTX)",
                data=st.session_state.dashboard_bytes,
                file_name=st.session_state.dashboard_filename or "SLEI_Dashboard.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            st.caption(
                "PDF export isn’t supported directly in Streamlit Cloud without a document-conversion service. "
                "If you want PDF output later, we can add a separate conversion step (e.g., external service or offline batch conversion)."
            )

            # Reset the flow for the next respondent (but keep download available until refresh)
            st.session_state.step = 1

else:
    st.session_state.step = 1
